from fastapi import FastAPI, File, UploadFile, HTTPException, Form
from fastapi.responses import JSONResponse
from fastapi.middleware.cors import CORSMiddleware
from typing import Optional
from uuid import uuid4
import os
import shutil
import json
import fitz  # PyMuPDF
from pptx import Presentation
from pathlib import Path

from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_community.embeddings import OpenAIEmbeddings
from langchain.chains import RetrievalQA
from langchain_community.chat_models import ChatOpenAI
from dotenv import load_dotenv

# Load .env if needed
load_dotenv()

app = FastAPI()

# CORS setup
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Directories
UPLOAD_DIR = "uploaded_docs"
VECTORSTORE_DIR = "vectorstores"
VECTORSTORE_PATH = os.path.join(VECTORSTORE_DIR, "global_store")
METADATA_FILE = os.path.join(VECTORSTORE_DIR, "metadata.json")
os.makedirs(UPLOAD_DIR, exist_ok=True)
os.makedirs(VECTORSTORE_DIR, exist_ok=True)

# Globals
global_db = None
metadata = {}

# Load existing vectorstore and metadata (if present)
try:
    if os.path.exists(os.path.join(VECTORSTORE_PATH, "index.faiss")):
        global_db = FAISS.load_local(VECTORSTORE_PATH, OpenAIEmbeddings()
                                     ,    allow_dangerous_deserialization=True
)
        print("✅ Loaded existing FAISS vectorstore.")

        # Check metadata presence
        test_docs = global_db.similarity_search("test", k=1)
        for doc in test_docs:
            if not doc.metadata.get("doc_id"):
                print("⚠️ FAISS documents may be missing metadata like 'doc_id'.")

    if os.path.exists(METADATA_FILE):
        with open(METADATA_FILE, "r") as f:
            metadata = json.load(f)
            print("✅ Loaded metadata.json with", len(metadata), "entries.")

except Exception as e:
    print(f"❌ Startup load failed: {e}")

@app.post("/upload")
async def upload_document(file: UploadFile = File(...)):
    if not (file.filename.endswith(".pdf") or file.filename.endswith(".pptx")):
        raise HTTPException(status_code=400, detail="Only PDF and PPTX files are supported.")

    doc_id = str(uuid4())
    ext = os.path.splitext(file.filename)[1].lower()
    file_path = os.path.join(UPLOAD_DIR, f"{doc_id}{ext}")

    # Save the uploaded file
    with open(file_path, "wb") as f:
        shutil.copyfileobj(file.file, f)

    # Extract text
    try:
        if ext == ".pdf":
            doc = fitz.open(file_path)
            text = "\n".join(page.get_text() for page in doc)
            doc.close()
        elif ext == ".pptx":
            prs = Presentation(file_path)
            text = "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Text extraction failed: {str(e)}")

    # Chunk + Embed
    try:
        splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
        docs = splitter.create_documents([text])

        # Tag each doc with its doc_id
        for d in docs:
            d.metadata = {"doc_id": doc_id}

        embeddings = OpenAIEmbeddings()
        new_db = FAISS.from_documents(docs, embeddings)

        global global_db
        if global_db is None:
            global_db = new_db
        else:
            global_db.merge_from(new_db)

        global_db.save_local(VECTORSTORE_PATH)

        # Save metadata
        metadata[doc_id] = file.filename
        with open(METADATA_FILE, "w") as f:
            json.dump(metadata, f)

        return {"message": "Document uploaded and indexed.", "doc_id": doc_id}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Embedding failed: {str(e)}")

@app.post("/ask")
async def ask_question(question: str = Form(...), doc_id: Optional[str] = Form(None)):
    try:
        if global_db is None:
            raise HTTPException(status_code=404, detail="No documents uploaded yet.")

        retriever = global_db.as_retriever()

        if doc_id and doc_id in metadata:
            results = global_db.similarity_search(question, k=10, filter={"doc_id": doc_id})

            context = "\n".join(doc.page_content for doc in results)

            if not context.strip():
                return {"answer": "Sorry, I couldn’t find a good answer in that document."}

            prompt = f"Answer based on this content:\n\n{context}\n\nQuestion: {question}"
            return {"answer": ChatOpenAI(model_name="gpt-4").predict(prompt)}

        # Default: search all docs
        qa_chain = RetrievalQA.from_chain_type(
            llm=ChatOpenAI(model_name="gpt-4"),
            retriever=retriever,
            chain_type="stuff"
        )
        result = qa_chain.run(question)
        return {"answer": result}

    except HTTPException as he:
        raise he
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error answering question: {str(e)}")

@app.get("/documents")
async def list_documents():
    return [{"doc_id": doc_id, "filename": name} for doc_id, name in metadata.items()]
